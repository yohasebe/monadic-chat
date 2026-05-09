#!/usr/bin/env python

import argparse
import json
import os
from docx import Document
import openpyxl
from pptx import Presentation

def extract_text_docx(doc_path):
    doc = Document(doc_path)
    parts = [para.text for para in doc.paragraphs]
    alt_texts = _docx_alt_texts(doc)
    for descr in alt_texts:
        parts.append(f"[Figure: {descr}]")
    return "\n".join(parts)


def _docx_alt_texts(doc):
    out = []
    try:
        body = doc.element.body
        for el in body.iter():
            tag = el.tag.split("}", 1)[-1] if "}" in el.tag else el.tag
            if tag != "docPr":
                continue
            label = (el.get("descr") or el.get("title") or "").strip()
            if label:
                out.append(label)
    except Exception:
        return out
    return out

def extract_text_xlsx(xlsx_path):
    workbook = openpyxl.load_workbook(xlsx_path)
    text = []
    for sheet in workbook:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value:
                    text.append(str(cell.value))
    return "\n".join(text)

def extract_text_pptx(ppt_path):
    prs = Presentation(ppt_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
            # Image alt text — preserves semantic context that would
            # otherwise vanish since we do not OCR images here.
            descr = _pptx_shape_descr(shape)
            if descr:
                text.append(f"[Figure: {descr}]")
        # Speaker notes are routinely the most information-dense surface
        # in a slide deck; surface them so the LLM can use them.
        notes = _pptx_notes(slide)
        if notes:
            text.append(f"[Speaker notes]\n{notes}")
    return "\n".join(text)


def _pptx_notes(slide):
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        nf = slide.notes_slide.notes_text_frame
        if nf is None:
            return ""
        return (nf.text or "").strip()
    except Exception:
        return ""


def _pptx_shape_descr(shape):
    try:
        el = shape.element
        for tag in ("nvSpPr", "nvPicPr", "nvGrpSpPr", "nvCxnSpPr"):
            ns = el.find(f".//{{*}}{tag}/{{*}}cNvPr")
            if ns is not None:
                descr = (ns.get("descr") or "").strip()
                if descr:
                    return descr
        return ""
    except Exception:
        return ""

def export_as_json(file_path):
    file_type = os.path.splitext(file_path)[1].lower()
    if file_type == '.docx':
        text = extract_text_docx(file_path)
    elif file_type == '.xlsx':
        text = extract_text_xlsx(file_path)
    elif file_type == '.pptx':
        text = extract_text_pptx(file_path)
    else:
        raise ValueError("Unsupported file type")

    data = {'text': text}
    json_data = json.dumps(data, ensure_ascii=False, indent=4)
    print(json_data)

    base_filename = os.path.splitext(os.path.basename(file_path))[0]
    output_filename = f"{base_filename}.json"
    with open(output_filename, 'w', encoding='utf-8') as f:
        f.write(json_data)

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Extract text from Office files and output as JSON.")
    parser.add_argument("file_path", help="Path to the Office file")
    args = parser.parse_args()

    if not os.path.exists(args.file_path):
        print(f"The specified file could not be found: {args.file_path}")
    else:
        export_as_json(args.file_path)
