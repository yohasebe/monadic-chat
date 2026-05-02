#!/usr/bin/env python3
"""Extract text + metadata from Microsoft Office files for Library import.

Supports:
  - .docx (python-docx): paragraphs preserved; Heading 1/2/3 styles
    become ATX headings so the Ruby importer can segment on them.
  - .xlsx (openpyxl): each sheet becomes an H1 section followed by a
    pipe-separated table of its non-empty rows.
  - .pptx (python-pptx): each slide becomes an H1 section using its
    title (or "Slide N" fallback), followed by the slide's body text.

Output is a JSON object on stdout:
{
  "title":       "<from doc properties or empty>",
  "author":      "<from doc properties or empty>",
  "format":      "docx" | "xlsx" | "pptx",
  "section_count": <int>,        # paragraphs / sheets / slides
  "markdown":    "<full content>"
}

Usage:
    python library_office_extractor.py /path/to/file.docx
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from typing import Iterable


# Style names that python-docx returns for built-in heading styles.
# We map them to ATX heading depth.
_DOCX_HEADING_STYLES = {
    "heading 1": 1, "title": 1,
    "heading 2": 2, "subtitle": 2,
    "heading 3": 3,
}


def extract_docx(path: str) -> dict:
    from docx import Document

    doc = Document(path)
    title = (doc.core_properties.title or "").strip()
    author = (doc.core_properties.author or "").strip()

    lines: list[str] = []
    para_count = 0
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            lines.append("")
            continue
        para_count += 1
        style = (para.style.name or "").strip().lower() if para.style else ""
        depth = _DOCX_HEADING_STYLES.get(style)
        if depth is not None:
            lines.append("#" * depth + " " + text)
        else:
            lines.append(text)
        lines.append("")

    # Tables — render as a single block per table, pipe-separated rows.
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            lines.append("\n".join(rows))
            lines.append("")

    # Image alt text — collected from inline + floating drawings. Authors
    # often place semantic context here (figure captions, screen-reader
    # descriptions) that would otherwise be lost since we do not OCR.
    alt_texts = _docx_alt_texts(doc)
    if alt_texts:
        lines.append("**Figures:** " + " / ".join(alt_texts))
        lines.append("")

    return {
        "title": title,
        "author": author,
        "format": "docx",
        "section_count": para_count,
        "markdown": "\n".join(lines).strip() + "\n",
    }


def _docx_alt_texts(doc) -> list[str]:
    out: list[str] = []
    try:
        body = doc.element.body
        # docPr (drawing properties) carries name/descr on inline + anchored shapes.
        for el in body.iter():
            tag = el.tag.split("}", 1)[-1] if "}" in el.tag else el.tag
            if tag != "docPr":
                continue
            descr = (el.get("descr") or "").strip()
            title_attr = (el.get("title") or "").strip()
            label = descr or title_attr
            if label:
                out.append(label)
    except Exception:
        return out
    return out


def extract_xlsx(path: str) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True, read_only=True)
    props = wb.properties
    title = (props.title or "").strip() if props else ""
    author = (props.creator or "").strip() if props else ""

    sections: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v).replace("\n", " ") for v in row]
            # Drop trailing empty cells so wide-but-sparse sheets stay readable.
            while cells and cells[-1] == "":
                cells.pop()
            if not cells:
                continue
            rows.append(" | ".join(cells))
        body = "\n".join(rows) if rows else "(empty sheet)"
        sections.append(f"# {sheet_name}\n\n{body}")

    wb.close()
    return {
        "title": title,
        "author": author,
        "format": "xlsx",
        "section_count": len(sections),
        "markdown": ("\n\n".join(sections) + "\n") if sections else "",
    }


def extract_pptx(path: str) -> dict:
    from pptx import Presentation

    prs = Presentation(path)
    props = prs.core_properties
    title = (props.title or "").strip() if props else ""
    author = (props.author or "").strip() if props else ""

    sections: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        body_parts: list[str] = []
        alt_texts: list[str] = []
        for shape in slide.shapes:
            # Image / picture shapes carry alt text on element.descr; capture
            # so figure context survives into the KB even though we do not
            # OCR the bitmap itself.
            descr = _shape_descr(shape)
            if descr:
                alt_texts.append(descr)
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if not text:
                    continue
                # Heuristic: the first non-empty line of the title placeholder
                # is the slide title. shape.is_placeholder + ph.type=='title'
                # would be more precise but adds version-specific imports.
                if not slide_title and shape == slide.shapes.title:
                    slide_title = text
                else:
                    body_parts.append(text)

        heading = slide_title or f"Slide {idx}"
        body_lines: list[str] = []
        if body_parts:
            body_lines.append("\n".join(body_parts))
        if alt_texts:
            body_lines.append("**Figures:** " + " / ".join(alt_texts))
        notes_text = _pptx_notes(slide)
        if notes_text:
            body_lines.append("**Speaker notes:**")
            body_lines.append(notes_text)
        body = "\n\n".join(body_lines) if body_lines else ""
        sections.append(f"# {heading}\n\n{body}".rstrip())

    return {
        "title": title,
        "author": author,
        "format": "pptx",
        "section_count": len(sections),
        "markdown": ("\n\n".join(sections) + "\n") if sections else "",
    }


def _pptx_notes(slide) -> str:
    """Best-effort speaker-notes extraction. Returns "" when absent."""
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return ""
        nf = notes_slide.notes_text_frame
        if nf is None:
            return ""
        text = (nf.text or "").strip()
        return text
    except Exception:
        return ""


def _shape_descr(shape) -> str:
    """Read an image/shape's alt text from the underlying XML element."""
    try:
        el = shape.element
        # nvSpPr / nvPicPr / nvGrpSpPr / nvCxnSpPr → cNvPr@descr
        for tag in ("nvSpPr", "nvPicPr", "nvGrpSpPr", "nvCxnSpPr"):
            ns = el.find(f".//{{*}}{tag}/{{*}}cNvPr")
            if ns is not None:
                descr = ns.get("descr") or ""
                if descr.strip():
                    return descr.strip()
        return ""
    except Exception:
        return ""


_DISPATCH = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


def extract(path: str) -> dict:
    ext = os.path.splitext(path)[1].lower()
    handler = _DISPATCH.get(ext)
    if handler is None:
        raise ValueError(f"Unsupported file extension: {ext}")
    return handler(path)


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("file_path", help="Path to the .docx / .xlsx / .pptx file")
    args = parser.parse_args(argv)

    try:
        result = extract(args.file_path)
    except FileNotFoundError:
        print(json.dumps({"error": "file_not_found", "path": args.file_path}), file=sys.stderr)
        return 2
    except ValueError as exc:
        print(json.dumps({"error": "unsupported_format", "message": str(exc)}), file=sys.stderr)
        return 3
    except Exception as exc:  # pragma: no cover - defensive
        print(json.dumps({"error": "extraction_failed", "message": str(exc)}), file=sys.stderr)
        return 1

    json.dump(result, sys.stdout, ensure_ascii=False)
    sys.stdout.write("\n")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
